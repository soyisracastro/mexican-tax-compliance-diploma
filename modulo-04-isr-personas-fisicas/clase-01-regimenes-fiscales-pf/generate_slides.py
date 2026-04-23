"""
Genera slides-clase-01.pptx desde slides-scripts.md.

Extrae solo el "Contenido Visual" de cada slide — los scripts del instructor
quedan fuera. Aplica el design system TodoConta (#0B5FFF · #0A1628 · #FAFAF7 · Inter).

Uso:
    python3 generate_slides.py
    python3 generate_slides.py ruta/slides-scripts.md salida.pptx

Requiere: python-pptx >= 1.0
    pip install python-pptx
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Design tokens (TodoConta) ────────────────────────────────────────────────
PRIMARY = RGBColor(0x0B, 0x5F, 0xFF)   # #0B5FFF — Azul Legal
BG      = RGBColor(0xFA, 0xFA, 0xF7)   # #FAFAF7 — Fondo cálido
DARK    = RGBColor(0x0A, 0x16, 0x28)   # #0A1628 — Texto / cover bg
MUTED   = RGBColor(0x55, 0x52, 0x49)   # #555249 — Texto secundario
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0xD4, 0xD1, 0xC9)   # #D4D1C9 — Separadores
ROW_ALT = RGBColor(0xF0, 0xF0, 0xEB)   # Filas alternas de tabla

FONT = "Inter"

# Dimensiones 16:9 (10" × 5.625")
W = Inches(10)
H = Inches(5.625)
M = Inches(0.52)  # margen horizontal


# ─── Helpers de shapes ───────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco


def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def _strip_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def _runs(p, text: str, size: float, base_color: RGBColor, bold_color: RGBColor = None):
    """Adds runs to paragraph with inline **bold** support."""
    if bold_color is None:
        bold_color = PRIMARY
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            r = p.add_run()
            r.text = part[2:-2]
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = bold_color
        elif part:
            r = p.add_run()
            r.text = part
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = False
            r.font.color.rgb = base_color


def _textbox(slide, left, top, width, height):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    txb.text_frame.word_wrap = True
    return txb


# ─── Parser ───────────────────────────────────────────────────────────────────

def parse(md_path: Path) -> dict:
    text = md_path.read_text(encoding="utf-8")
    lines = text.split("\n")

    result = {"title": "", "subtitle": "", "meta": [], "slides": []}
    i = 0

    # Cabecera del deck
    while i < len(lines):
        line = lines[i]
        if line.startswith("# ") and "BLOQUE" not in line:
            result["title"] = line[2:].strip()
        elif line.startswith("## ") and "SLIDE" not in line and "###" not in line:
            result["subtitle"] = line[3:].strip()
        elif re.match(r"^\*\*Duración\*\*", line) or re.match(r"^\*\*Actualización", line) or re.match(r"^\*\*Contenido", line):
            result["meta"].append(re.sub(r"\*\*", "", line).strip())
        elif re.match(r"^## SLIDE|^# BLOQUE", line):
            break
        i += 1

    # Slides
    current = None
    in_visual = False
    visual_lines: list[str] = []

    def flush():
        nonlocal current, visual_lines
        if current is not None:
            current["content"] = _trim(visual_lines)
            result["slides"].append(current)
        current = None
        visual_lines = []

    while i < len(lines):
        line = lines[i]

        if re.match(r"^# BLOQUE", line):
            flush()
            title = re.sub(r"\s*\{icon:[^}]+\}", "", line[2:]).strip()
            result["slides"].append({"type": "block", "title": title, "content": []})
            in_visual = False

        elif re.match(r"^## SLIDE \d+:", line):
            flush()
            slide_title = re.sub(r"^## SLIDE \d+:\s*", "", line).strip()
            current = {"type": "content", "title": slide_title, "content": []}
            in_visual = False

        elif line.startswith("### Contenido Visual"):
            in_visual = True

        elif line.startswith("### Script"):
            in_visual = False

        elif line.strip() == "---":
            pass

        elif in_visual and current is not None:
            visual_lines.append(line)

        i += 1

    flush()
    return result


def _trim(lines: list[str]) -> list[str]:
    while lines and lines[0].strip() == "":
        lines.pop(0)
    while lines and lines[-1].strip() == "":
        lines.pop()
    return lines


def _parse_table(lines: list[str]):
    headers, rows = [], []
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|"):
            continue
        if re.match(r"^\|[-:\s|]+\|$", stripped):
            continue  # separador
        cells = [c.strip() for c in stripped.strip("|").split("|")]
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


# ─── Slide builders ───────────────────────────────────────────────────────────

def build_cover(prs: Presentation, title: str, subtitle: str, meta: list[str]):
    slide = _blank_slide(prs)
    _bg(slide, DARK)

    # Barra de color izquierda
    _rect(slide, M, Inches(1.4), Inches(0.07), Inches(2.8), PRIMARY)

    # Título
    txb = _textbox(slide, Inches(0.78), Inches(1.3), Inches(8.7), Inches(1.8))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    _runs(p, title, 30, WHITE, WHITE)
    p.runs[0].font.bold = True

    # Subtítulo
    txb2 = _textbox(slide, Inches(0.78), Inches(3.15), Inches(8.7), Inches(0.65))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    _runs(p2, subtitle, 17, PRIMARY, PRIMARY)

    # Meta
    if meta:
        txb3 = _textbox(slide, Inches(0.78), Inches(3.9), Inches(8.7), Inches(1.3))
        tf3 = txb3.text_frame
        first = True
        for line in meta:
            p3 = tf3.paragraphs[0] if first else tf3.add_paragraph()
            first = False
            _runs(p3, line, 12, MUTED, MUTED)

    # Barra inferior azul
    _rect(slide, 0, H - Inches(0.2), W, Inches(0.2), PRIMARY)


def build_block(prs: Presentation, title: str):
    slide = _blank_slide(prs)
    _bg(slide, PRIMARY)

    txb = _textbox(slide, M, Inches(1.7), W - 2 * M, Inches(2.3))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Línea decorativa
    _rect(slide, Inches(3.5), Inches(4.0), Inches(3.0), Inches(0.05), WHITE)


def build_content(prs: Presentation, title: str, content: list[str], num: int):
    slide = _blank_slide(prs)
    _bg(slide, BG)

    # Acento azul junto al título
    _rect(slide, M, Inches(0.32), Inches(0.07), Inches(0.65), PRIMARY)

    # Título
    txb_t = _textbox(slide, Inches(0.74), Inches(0.28), Inches(8.8), Inches(0.75))
    tf_t = txb_t.text_frame
    p_t = tf_t.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = _strip_bold(title)
    r_t.font.name = FONT
    r_t.font.size = Pt(21)
    r_t.font.bold = True
    r_t.font.color.rgb = DARK

    # Línea separadora
    _rect(slide, M, Inches(1.08), W - 2 * M, Inches(0.025), BORDER)

    # Número de slide
    txb_n = _textbox(slide, W - Inches(1.0), H - Inches(0.38), Inches(0.82), Inches(0.28))
    tf_n = txb_n.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.alignment = PP_ALIGN.RIGHT
    r_n = p_n.add_run()
    r_n.text = str(num)
    r_n.font.name = FONT
    r_n.font.size = Pt(9)
    r_n.font.color.rgb = MUTED

    # Área de contenido
    C_LEFT = M
    C_TOP  = Inches(1.2)
    C_W    = W - 2 * M
    C_H    = H - C_TOP - Inches(0.45)

    # Detectar tabla
    table_idx = next((i for i, l in enumerate(content) if l.strip().startswith("|")), None)

    if table_idx is not None:
        pre  = _trim(content[:table_idx])
        tbl_lines = []
        post = []
        in_tbl = True
        for line in content[table_idx:]:
            if in_tbl and line.strip().startswith("|"):
                tbl_lines.append(line)
            else:
                in_tbl = False
                post.append(line)
        post = _trim(post)

        y = C_TOP

        # Texto previo a la tabla
        if pre:
            pre_h = Inches(0.38 * len([l for l in pre if l.strip()]) + 0.1)
            pre_h = min(pre_h, Inches(1.2))
            txb_pre = _textbox(slide, C_LEFT, y, C_W, pre_h)
            tf_pre = txb_pre.text_frame
            first = True
            for line in pre:
                if line.strip() == "":
                    continue
                p = tf_pre.paragraphs[0] if first else tf_pre.add_paragraph()
                first = False
                p.space_before = Pt(3)
                is_bullet = line.strip().startswith("-")
                text = line.strip().lstrip("- ").strip() if is_bullet else line.strip()
                if is_bullet:
                    _runs(p, "• " + text, 14, DARK)
                else:
                    _runs(p, text, 14, DARK)
            y += pre_h + Inches(0.08)

        # Tabla
        headers, rows = _parse_table(tbl_lines)
        if headers and len(headers) > 0:
            n_rows = 1 + len(rows)
            row_h  = Inches(0.34)
            tbl_h  = row_h * n_rows
            remaining = C_TOP + C_H - y - Inches(0.1)
            tbl_h = min(tbl_h, remaining)

            tbl_shape = slide.shapes.add_table(n_rows, len(headers), C_LEFT, y, C_W, tbl_h)
            tbl = tbl_shape.table

            # Ancho de columnas
            col_w = C_W // len(headers)
            for ci in range(len(headers)):
                tbl.columns[ci].width = col_w

            # Cabecera
            for ci, hdr in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = _strip_bold(hdr)
                r.font.name = FONT
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = WHITE

            # Filas de datos
            for ri, row in enumerate(rows):
                bg = WHITE if ri % 2 == 0 else ROW_ALT
                for ci in range(len(headers)):
                    cell = tbl.cell(ri + 1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg
                    val = row[ci] if ci < len(row) else ""
                    is_bold_cell = "**" in val
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
                    r = p.add_run()
                    r.text = _strip_bold(val)
                    r.font.name = FONT
                    r.font.size = Pt(11)
                    r.font.bold = is_bold_cell
                    r.font.color.rgb = PRIMARY if is_bold_cell else DARK

            y += tbl_h + Inches(0.1)

        # Texto posterior
        if post:
            remaining = C_TOP + C_H - y
            if remaining > Inches(0.25):
                txb_post = _textbox(slide, C_LEFT, y, C_W, remaining)
                tf_post = txb_post.text_frame
                first = True
                for line in post:
                    if line.strip() == "":
                        continue
                    p = tf_post.paragraphs[0] if first else tf_post.add_paragraph()
                    first = False
                    p.space_before = Pt(3)
                    is_bullet = line.strip().startswith("-")
                    text = line.strip().lstrip("- ").strip() if is_bullet else line.strip()
                    _runs(p, ("• " if is_bullet else "") + text, 13, MUTED)

    else:
        # Solo texto / bullets
        txb_c = _textbox(slide, C_LEFT, C_TOP, C_W, C_H)
        tf_c = txb_c.text_frame
        first = True
        for line in content:
            if line.strip() == "":
                if not first:
                    p = tf_c.add_paragraph()
                    p.space_before = Pt(2)
                continue
            p = tf_c.paragraphs[0] if first else tf_c.add_paragraph()
            first = False
            p.space_before = Pt(7)

            is_bullet = line.strip().startswith("-")
            text = line.strip().lstrip("- ").strip() if is_bullet else line.strip()
            _runs(p, ("• " if is_bullet else "") + text, 15, DARK)


# ─── Main ─────────────────────────────────────────────────────────────────────

def generate(md_path: Path, out_path: Path):
    data = parse(md_path)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Cover
    build_cover(prs, data["title"], data["subtitle"], data["meta"])

    slide_num = 1
    for s in data["slides"]:
        if s["type"] == "block":
            build_block(prs, s["title"])
        elif s["type"] == "content":
            slide_num += 1
            build_content(prs, s["title"], s["content"], slide_num)

    prs.save(out_path)
    print(f"✓ {out_path.name}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    here = Path(__file__).parent
    md   = Path(sys.argv[1]) if len(sys.argv) > 1 else here / "slides-scripts.md"
    out  = Path(sys.argv[2]) if len(sys.argv) > 2 else here / "slides-clase-01-modulo-04.pptx"
    generate(md, out)
